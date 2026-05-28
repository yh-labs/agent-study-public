"""
PRD 기반 반려동물 건강 기록 웹앱 발표 자료 생성기
python-pptx 사용
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import datetime

# ───────────────────────────────────────────────
# 색상 팔레트 (Teal Trust — 헬스/웰니스 테마)
# ───────────────────────────────────────────────
C_PRIMARY   = RGBColor(0x02, 0x80, 0x90)   # teal #028090
C_ACCENT    = RGBColor(0x02, 0xC3, 0x9A)   # mint #02C39A
C_DARK      = RGBColor(0x1E, 0x29, 0x3B)   # slate #1E293B
C_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT_BG  = RGBColor(0xF0, 0xFA, 0xFB)   # 연한 민트 흰색
C_MUTED     = RGBColor(0x64, 0x74, 0x8B)   # 회색
C_CARD_BG   = RGBColor(0xFF, 0xFF, 0xFF)
C_HEADER_BG = RGBColor(0x02, 0x80, 0x90)
C_TAG_BG    = RGBColor(0xE0, 0xF5, 0xF7)   # 연 teal
C_TAG_TEXT  = RGBColor(0x02, 0x60, 0x70)
C_MUST      = RGBColor(0x02, 0x80, 0x90)
C_SHOULD    = RGBColor(0x02, 0xC3, 0x9A)
C_COULD     = RGBColor(0x94, 0xA3, 0xB8)
C_ROW_ALT   = RGBColor(0xF0, 0xFA, 0xFB)

# 슬라이드 사이즈 (LAYOUT_16x9): 10" x 5.625"
W = Inches(10)
H = Inches(5.625)


def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


# ───────────────────────────────────────────────
# 헬퍼 함수들
# ───────────────────────────────────────────────

def blank_slide(prs):
    blank_layout = prs.slide_layouts[6]
    return prs.slides.add_slide(blank_layout)


def bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def rect(slide, x, y, w, h, fill_color, line_color=None, line_width=0):
    from pptx.util import Pt as _Pt
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE = 1
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = _Pt(line_width)
    else:
        shape.line.fill.background()
    return shape


def txt(slide, text, x, y, w, h,
        size=18, bold=False, color=None, align=PP_ALIGN.LEFT,
        italic=False, wrap=True, margin_left=0.05):
    txBox = slide.shapes.add_textbox(
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color if color else C_DARK
    txBox.left = Inches(x)
    return txBox


def multiline_txt(slide, lines, x, y, w, h,
                  size=14, color=None, bold=False, line_spacing_pt=6,
                  align=PP_ALIGN.LEFT):
    """
    lines: list of (text, bold, color) tuples  or  plain strings
    """
    txBox = slide.shapes.add_textbox(
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if isinstance(line, str):
            t, b, c = line, bold, (color or C_DARK)
        else:
            t, b, c = line
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(line_spacing_pt)
        run = p.add_run()
        run.text = t
        run.font.size = Pt(size)
        run.font.bold = b
        run.font.color.rgb = c
    return txBox


def section_header(slide, title, y_pos=0.0, h=0.55):
    rect(slide, 0, y_pos, 10, h, C_PRIMARY)
    txt(slide, title,
        x=0.35, y=y_pos + 0.05, w=9, h=h - 0.05,
        size=20, bold=True, color=C_WHITE)


def card(slide, x, y, w, h, title, body_lines,
         title_size=14, body_size=12.5, accent_bar=True):
    rect(slide, x, y, w, h, C_CARD_BG)
    if accent_bar:
        rect(slide, x, y, 0.07, h, C_ACCENT)
    txt(slide, title,
        x=x + 0.18, y=y + 0.07, w=w - 0.25, h=0.3,
        size=title_size, bold=True, color=C_PRIMARY)
    lines = []
    for line in body_lines:
        if isinstance(line, str):
            lines.append((line, False, C_DARK))
        else:
            lines.append(line)
    multiline_txt(slide, lines,
                  x=x + 0.18, y=y + 0.4, w=w - 0.25, h=h - 0.5,
                  size=body_size, line_spacing_pt=3)


def stat_box(slide, x, y, w, h, value, label, val_color=None):
    rect(slide, x, y, w, h, C_CARD_BG)
    rect(slide, x, y, w, 0.06, C_ACCENT)  # top accent line
    txt(slide, value,
        x=x + 0.1, y=y + 0.15, w=w - 0.2, h=0.55,
        size=32, bold=True, color=val_color or C_PRIMARY,
        align=PP_ALIGN.CENTER)
    txt(slide, label,
        x=x + 0.1, y=y + 0.7, w=w - 0.2, h=0.35,
        size=11, color=C_MUTED, align=PP_ALIGN.CENTER)


def pill_tag(slide, x, y, w, h, text, bg_c=None, text_c=None):
    rect(slide, x, y, w, h, bg_c or C_TAG_BG)
    txt(slide, text,
        x=x, y=y + 0.01, w=w, h=h,
        size=11, bold=True, color=text_c or C_TAG_TEXT,
        align=PP_ALIGN.CENTER)


# ───────────────────────────────────────────────
# 슬라이드 생성 함수들
# ───────────────────────────────────────────────

def slide_01_title(prs):
    """제목 슬라이드"""
    s = blank_slide(prs)
    bg(s, C_PRIMARY)

    # 왼쪽 민트 세로 바
    rect(s, 0, 0, 0.5, 5.625, C_ACCENT)

    # 상단 장식 원
    from pptx.util import Inches as I2
    circ = s.shapes.add_shape(9, I2(7.8), I2(-0.5), I2(3), I2(3))  # OVAL
    circ.fill.solid()
    circ.fill.fore_color.rgb = RGBColor(0x00, 0x6B, 0x78)
    circ.line.fill.background()

    # 하단 장식 원
    circ2 = s.shapes.add_shape(9, I2(8.5), I2(3.8), I2(2.5), I2(2.5))
    circ2.fill.solid()
    circ2.fill.fore_color.rgb = RGBColor(0x02, 0xA0, 0xB4)
    circ2.line.fill.background()

    # 서브타이틀
    txt(s, "Product Requirements Document",
        x=0.8, y=1.35, w=7.5, h=0.4,
        size=13, color=C_ACCENT, bold=False, italic=True)

    # 메인 타이틀
    txt(s, "반려동물 건강 기록",
        x=0.8, y=1.8, w=8, h=0.75,
        size=40, bold=True, color=C_WHITE)
    txt(s, "웹앱 서비스",
        x=0.8, y=2.5, w=8, h=0.65,
        size=40, bold=True, color=C_WHITE)

    # 구분선
    line = s.shapes.add_shape(20, Inches(0.8), Inches(3.25), Inches(5), Emu(0))
    line.line.color.rgb = C_ACCENT
    line.line.width = Pt(2)

    # 날짜/기간
    txt(s, "2026.05.19 ~ 2026.07.22   |   9주 프로젝트",
        x=0.8, y=3.45, w=8, h=0.35,
        size=13, color=RGBColor(0xB2, 0xE4, 0xEA))

    # 기술 스택 태그
    tags = ["React + Vite", "Express.js", "PostgreSQL", "Supabase", "Vercel"]
    for i, tag in enumerate(tags):
        pill_tag(s, 0.8 + i * 1.65, 4.0, 1.5, 0.35, tag,
                 bg_c=RGBColor(0x01, 0x60, 0x70),
                 text_c=RGBColor(0xB2, 0xE4, 0xEA))

    return s


def slide_02_overview(prs):
    """프로젝트 개요"""
    s = blank_slide(prs)
    bg(s, C_LIGHT_BG)
    section_header(s, "프로젝트 개요", y_pos=0.0)

    # 왼쪽: 핵심 설명
    card(s, 0.3, 0.7, 5.0, 1.6,
         "서비스 한 줄 요약",
         ["강아지·고양이 보호자가 반려동물의 데일리 체중·음수량을",
          "간편하게 기록하고, 월간 건강 리포트로 장기 추이를 확인하는",
          "PWA 웹앱 서비스"],
         title_size=15, body_size=13)

    # 오른쪽: 팀/기간 정보
    info = [
        ("프로젝트명", "반려동물 건강 기록 웹앱"),
        ("기간",      "9주 (2026-05-19 ~ 2026-07-22)"),
        ("팀 구성",   "PM 1 · 기획 1 · 디자인 1 · 개발 2"),
        ("기술 스택", "React/Vite · Express.js · PostgreSQL"),
        ("배포",      "Vercel (프론트 + API 통합)"),
        ("데이터베이스", "Supabase (PostgreSQL + Storage)"),
    ]

    row_y = 0.7
    for label, value in info:
        rect(s, 5.6, row_y, 1.55, 0.42, C_PRIMARY)
        txt(s, label,
            x=5.65, y=row_y + 0.03, w=1.45, h=0.38,
            size=12, bold=True, color=C_WHITE)
        rect(s, 7.15, row_y, 2.55, 0.42, C_CARD_BG)
        txt(s, value,
            x=7.2, y=row_y + 0.03, w=2.45, h=0.38,
            size=12, color=C_DARK)
        row_y += 0.44

    # 하단: 목표 3가지
    goals = [
        ("1분 기록", "매일 1분 이내\n체중·음수량 입력"),
        ("자동 분석", "월간 리포트로\n이상 징후 인지"),
        ("병원 공유", "동물병원 방문 시\n건강 이력 제공"),
    ]
    for i, (title, desc) in enumerate(goals):
        bx = 0.3 + i * 3.2
        rect(s, bx, 4.5, 3.0, 0.9, C_PRIMARY)
        txt(s, title,
            x=bx + 0.1, y=4.52, w=2.8, h=0.35,
            size=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        txt(s, desc.replace("\n", "  "),
            x=bx + 0.1, y=4.87, w=2.8, h=0.45,
            size=11, color=RGBColor(0xB2, 0xE4, 0xEA), align=PP_ALIGN.CENTER)

    return s


def slide_03_problem(prs):
    """문제 정의 & 목표"""
    s = blank_slide(prs)
    bg(s, C_LIGHT_BG)
    section_header(s, "배경 및 문제 정의", y_pos=0.0)

    # 문제
    txt(s, "문제",
        x=0.3, y=0.7, w=4.5, h=0.35,
        size=16, bold=True, color=C_PRIMARY)
    problems = [
        "반려동물 체중 급변·수분 섭취 이상은 질병 초기 징후",
        "→ 대부분의 보호자가 수기 관리하거나 아예 기록하지 않음",
        "동물병원 방문 시 최근 이력 제시 불가 → 진단 어려움",
    ]
    for i, p in enumerate(problems):
        rect(s, 0.3, 1.1 + i * 0.62, 4.6, 0.5, C_CARD_BG)
        rect(s, 0.3, 1.1 + i * 0.62, 0.06, 0.5,
             C_MUST if i == 0 else C_MUTED)
        txt(s, p,
            x=0.5, y=1.13 + i * 0.62, w=4.3, h=0.42,
            size=12.5, color=C_DARK)

    # 목표
    txt(s, "목표",
        x=5.3, y=0.7, w=4.5, h=0.35,
        size=16, bold=True, color=C_PRIMARY)
    goals = [
        "매일 1분 이내 체중·음수량 기록 가능",
        "월간 리포트로 보호자 스스로 이상 징후 인지",
        "동물병원 방문 시 공유 가능한 건강 이력 제공",
    ]
    for i, g in enumerate(goals):
        rect(s, 5.3, 1.1 + i * 0.62, 4.4, 0.5, C_CARD_BG)
        rect(s, 5.3, 1.1 + i * 0.62, 0.06, 0.5, C_ACCENT)
        txt(s, g,
            x=5.5, y=1.13 + i * 0.62, w=4.1, h=0.42,
            size=12.5, color=C_DARK)

    # 비기능 요구사항 요약 바
    rect(s, 0.3, 3.25, 9.4, 1.05, C_PRIMARY)
    nfr = [
        ("플랫폼", "PWA / Chrome 90+\nSafari 14+ / Firefox 90+"),
        ("응답속도", "기록 저장 ≤ 1초"),
        ("오프라인", "Service Worker\n자동 동기화"),
        ("보안", "건강 데이터\n암호화 저장"),
        ("데이터보존", "탈퇴 후 30일\n보관 후 삭제"),
    ]
    for i, (label, val) in enumerate(nfr):
        bx = 0.45 + i * 1.85
        txt(s, label,
            x=bx, y=3.28, w=1.75, h=0.3,
            size=11, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)
        txt(s, val,
            x=bx, y=3.57, w=1.75, h=0.65,
            size=11, color=C_WHITE, align=PP_ALIGN.CENTER)

    return s


def slide_04_users(prs):
    """대상 사용자"""
    s = blank_slide(prs)
    bg(s, C_LIGHT_BG)
    section_header(s, "대상 사용자", y_pos=0.0)

    segments = [
        ("1차 타겟",
         "강아지·고양이를 키우는\n20–40대 보호자",
         "핵심 타겟 — 앱 친숙도 높고\n일상 건강 관리 수요 있음",
         C_PRIMARY),
        ("2차 타겟",
         "노령견·노령묘\n보호자",
         "건강 관리 수요 매우 높음\n정기적 체크 필요",
         C_ACCENT),
        ("3차 타겟",
         "다두 가정\n보호자",
         "멀티 펫 지원 니즈\n각 반려동물 개별 관리",
         C_MUTED),
    ]

    for i, (seg, desc, detail, color) in enumerate(segments):
        bx = 0.35 + i * 3.15
        rect(s, bx, 0.7, 2.9, 4.3, C_CARD_BG)
        rect(s, bx, 0.7, 2.9, 0.55, color)
        txt(s, seg,
            x=bx + 0.1, y=0.73, w=2.7, h=0.45,
            size=17, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        txt(s, desc,
            x=bx + 0.1, y=1.4, w=2.7, h=1.0,
            size=15, bold=True, color=C_DARK, align=PP_ALIGN.CENTER)
        # 구분선
        line = s.shapes.add_shape(20,
                                   Inches(bx + 0.3), Inches(2.55),
                                   Inches(2.3), Emu(0))
        line.line.color.rgb = C_TAG_BG
        line.line.width = Pt(1.5)
        txt(s, detail,
            x=bx + 0.1, y=2.7, w=2.7, h=1.2,
            size=12, color=C_MUTED, align=PP_ALIGN.CENTER)

    # KPI 미리보기
    rect(s, 0.3, 5.0, 9.4, 0.45, RGBColor(0xE0, 0xF5, 0xF7))
    txt(s, "런칭 3개월 목표   MAU 5,000   7일 리텐션 ≥ 40%   데일리 기록 완료율 ≥ 60%   리포트 열람율 ≥ 70%",
        x=0.4, y=5.02, w=9.2, h=0.38,
        size=12, color=C_TAG_TEXT, align=PP_ALIGN.CENTER)

    return s


def slide_05_feature_profile_weight(prs):
    """핵심 기능 (1) 프로필 관리 + 체중 기록"""
    s = blank_slide(prs)
    bg(s, C_LIGHT_BG)
    section_header(s, "핵심 기능 (1) — 반려동물 프로필 관리 + 데일리 체중 기록")

    # 프로필 관리 카드
    rect(s, 0.3, 0.68, 4.6, 4.5, C_CARD_BG)
    rect(s, 0.3, 0.68, 4.6, 0.45, C_PRIMARY)
    txt(s, "반려동물 프로필 관리",
        x=0.45, y=0.7, w=4.3, h=0.4,
        size=15, bold=True, color=C_WHITE)

    profile_items = [
        "이름, 종, 품종, 생년월일, 성별, 중성화 여부 등록",
        "프로필 사진 업로드",
        "멀티 펫 지원 (다두 가정 대응)",
    ]
    for i, item in enumerate(profile_items):
        rect(s, 0.45, 1.25 + i * 0.72, 4.3, 0.58, C_LIGHT_BG)
        txt(s, f"  ·  {item}",
            x=0.5, y=1.28 + i * 0.72, w=4.2, h=0.52,
            size=13, color=C_DARK)

    # 미니 차트 시각화 (모의)
    rect(s, 0.45, 3.45, 4.3, 1.55, C_LIGHT_BG)
    txt(s, "반려동물 등록 예시",
        x=0.55, y=3.48, w=4.1, h=0.3,
        size=11, bold=True, color=C_MUTED)
    fields = [("이름", "코코"), ("종", "강아지"), ("품종", "포메라니안"),
              ("생년월일", "2021-03-15"), ("체중(최근)", "3.2 kg")]
    for j, (k, v) in enumerate(fields):
        cx = 0.55 + (j % 2) * 2.1
        cy = 3.82 + (j // 2) * 0.42
        rect(s, cx, cy, 0.8, 0.32, C_TAG_BG)
        txt(s, k, x=cx, y=cy + 0.02, w=0.8, h=0.28,
            size=10, bold=True, color=C_TAG_TEXT, align=PP_ALIGN.CENTER)
        txt(s, v, x=cx + 0.85, y=cy + 0.02, w=1.1, h=0.28,
            size=10, color=C_DARK)

    # 체중 기록 카드
    rect(s, 5.1, 0.68, 4.6, 4.5, C_CARD_BG)
    rect(s, 5.1, 0.68, 4.6, 0.45, C_ACCENT)
    txt(s, "데일리 체중 기록",
        x=5.25, y=0.7, w=4.3, h=0.4,
        size=15, bold=True, color=C_WHITE)

    weight_items = [
        "날짜별 체중 입력 (kg, 소수점 1자리)",
        "최근 7일 추이 미니 차트 홈 노출",
        "전일 대비 변화량 표시 (예: −0.2 kg ▼)",
        "권장 체중 범위 표시 (품종·나이 기반)",
    ]
    for i, item in enumerate(weight_items):
        rect(s, 5.25, 1.25 + i * 0.55, 4.3, 0.45, C_LIGHT_BG)
        txt(s, f"  ·  {item}",
            x=5.3, y=1.27 + i * 0.55, w=4.2, h=0.41,
            size=12.5, color=C_DARK)

    # 체중 차트 시뮬레이션
    rect(s, 5.25, 3.52, 4.3, 1.45, C_LIGHT_BG)
    txt(s, "최근 7일 체중 추이 (예시)",
        x=5.35, y=3.55, w=4.1, h=0.3,
        size=11, bold=True, color=C_MUTED)

    weights = [3.4, 3.3, 3.35, 3.2, 3.25, 3.18, 3.2]
    days    = ["화", "수", "목", "금", "토", "일", "월"]
    bar_w   = 0.42
    max_w   = max(weights)
    min_w   = min(weights)
    for j, (day, w_val) in enumerate(zip(days, weights)):
        bx = 5.35 + j * 0.56
        norm = (w_val - min_w + 0.1) / (max_w - min_w + 0.2)
        bar_h = norm * 0.65
        bar_y = 4.5 - bar_h
        bar_color = C_PRIMARY if j < 6 else C_ACCENT
        rect(s, bx, bar_y, bar_w, bar_h, bar_color)
        txt(s, day, x=bx, y=4.52, w=bar_w, h=0.22,
            size=10, color=C_MUTED, align=PP_ALIGN.CENTER)
        txt(s, f"{w_val}", x=bx, y=bar_y - 0.22, w=bar_w, h=0.22,
            size=9, color=C_PRIMARY, align=PP_ALIGN.CENTER, bold=True)

    return s


def slide_06_feature_water_report(prs):
    """핵심 기능 (2) 음수량 + 월간 리포트"""
    s = blank_slide(prs)
    bg(s, C_LIGHT_BG)
    section_header(s, "핵심 기능 (2) — 데일리 음수량 기록 + 월간 건강 리포트")

    # 음수량 카드
    rect(s, 0.3, 0.68, 4.6, 4.5, C_CARD_BG)
    rect(s, 0.3, 0.68, 4.6, 0.45, C_PRIMARY)
    txt(s, "데일리 음수량 기록",
        x=0.45, y=0.7, w=4.3, h=0.4,
        size=15, bold=True, color=C_WHITE)

    water_items = [
        "날짜별 음수량 입력 (ml 단위)",
        "물그릇 용량 설정 → 잔여량 역산 지원",
        "권장 음수량 범위 표시 (체중 기반)",
    ]
    for i, item in enumerate(water_items):
        rect(s, 0.45, 1.25 + i * 0.58, 4.3, 0.47, C_LIGHT_BG)
        txt(s, f"  ·  {item}",
            x=0.5, y=1.28 + i * 0.58, w=4.2, h=0.42,
            size=13, color=C_DARK)

    # 음수량 시각화
    rect(s, 0.45, 3.1, 4.3, 1.85, C_LIGHT_BG)
    txt(s, "권장 음수량 계산 예시 (체중 기반)",
        x=0.55, y=3.13, w=4.1, h=0.28,
        size=11, bold=True, color=C_MUTED)
    txt(s, "체중 3.2 kg",
        x=0.55, y=3.45, w=4.1, h=0.3,
        size=13, bold=True, color=C_DARK, align=PP_ALIGN.CENTER)

    # 게이지 바
    rect(s, 0.7, 3.85, 3.9, 0.35, RGBColor(0xE2, 0xE8, 0xF0))
    rect(s, 0.7, 3.85, 2.73, 0.35, C_ACCENT)  # 70% 기록
    txt(s, "오늘 음수량: 224 ml / 권장 320 ml (70%)",
        x=0.55, y=4.25, w=4.1, h=0.3,
        size=11, color=C_MUTED, align=PP_ALIGN.CENTER)

    # 월간 리포트 카드
    rect(s, 5.1, 0.68, 4.6, 4.5, C_CARD_BG)
    rect(s, 5.1, 0.68, 4.6, 0.45, C_ACCENT)
    txt(s, "월간 건강 리포트",
        x=5.25, y=0.7, w=4.3, h=0.4,
        size=15, bold=True, color=C_WHITE)

    report_items = [
        "매월 1일 자동 생성 + 푸시 알림",
        "월간 평균 체중·음수량 및 전월 대비 변화",
        "체중·음수량 라인 차트 (30일 추이)",
        "이상 징후 강조 표시 (임계값 초과)",
        "PDF 내보내기 및 공유 기능",
    ]
    for i, item in enumerate(report_items):
        rect(s, 5.25, 1.25 + i * 0.52, 4.3, 0.42, C_LIGHT_BG)
        txt(s, f"  ·  {item}",
            x=5.3, y=1.27 + i * 0.52, w=4.2, h=0.38,
            size=12.5, color=C_DARK)

    # 리포트 요약 박스
    rect(s, 5.25, 3.95, 4.3, 1.0, C_PRIMARY)
    txt(s, "5월 건강 요약 (예시)",
        x=5.35, y=3.98, w=4.1, h=0.28,
        size=11, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)
    report_stats = [
        ("평균 체중", "3.24 kg", "전월 대비 −0.05 kg"),
        ("평균 음수량", "298 ml", "권장 320 ml 대비 93%"),
    ]
    for j, (label, value, note) in enumerate(report_stats):
        bx = 5.35 + j * 2.15
        txt(s, label, x=bx, y=4.3, w=2.0, h=0.22,
            size=10, color=C_ACCENT, align=PP_ALIGN.CENTER)
        txt(s, value, x=bx, y=4.52, w=2.0, h=0.3,
            size=16, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        txt(s, note, x=bx, y=4.82, w=2.0, h=0.2,
            size=9, color=RGBColor(0xB2, 0xE4, 0xEA), align=PP_ALIGN.CENTER)

    return s


def slide_07_feature_alerts(prs):
    """핵심 기능 (3) 알림 + 사용자 흐름"""
    s = blank_slide(prs)
    bg(s, C_LIGHT_BG)
    section_header(s, "핵심 기능 (3) — 알림 + 사용자 흐름")

    # 알림 종류 카드 2개
    alert_cards = [
        ("기록 리마인더 알림",
         ["매일 정해진 시간에 기록 리마인더 발송",
          "보호자가 시간대 직접 설정",
          "PWA Push Notification 활용"],
         C_PRIMARY),
        ("이상 징후 즉시 알림",
         ["전일 대비 체중 10% 이상 변화 감지",
          "설정 임계값 초과 시 즉시 알림",
          "월간 리포트 자동 생성 알림 포함"],
         C_ACCENT),
    ]
    for i, (title, items, color) in enumerate(alert_cards):
        bx = 0.3 + i * 4.85
        rect(s, bx, 0.68, 4.5, 3.0, C_CARD_BG)
        rect(s, bx, 0.68, 4.5, 0.45, color)
        txt(s, title, x=bx + 0.15, y=0.7, w=4.2, h=0.4,
            size=14, bold=True, color=C_WHITE)
        for j, item in enumerate(items):
            rect(s, bx + 0.15, 1.25 + j * 0.62, 4.2, 0.5, C_LIGHT_BG)
            txt(s, f"  ·  {item}",
                x=bx + 0.2, y=1.28 + j * 0.62, w=4.0, h=0.44,
                size=12.5, color=C_DARK)

    # 사용자 흐름
    rect(s, 0.3, 3.85, 9.4, 0.35, C_PRIMARY)
    txt(s, "사용자 흐름 (User Flow)",
        x=0.45, y=3.87, w=9.1, h=0.3,
        size=14, bold=True, color=C_WHITE)

    flow_steps = [
        "온보딩",
        "프로필 등록",
        "홈 화면",
        "[+] 오늘 기록",
        "체중 입력",
        "음수량 입력",
        "저장 완료",
        "월간 리포트",
    ]
    step_colors = [C_PRIMARY, C_PRIMARY, C_ACCENT, C_ACCENT,
                   C_ACCENT, C_ACCENT, C_ACCENT, C_PRIMARY]
    arrow_color = C_MUTED
    sw = 1.08
    for i, (step, sc) in enumerate(zip(flow_steps, step_colors)):
        bx = 0.35 + i * 1.17
        rect(s, bx, 4.28, sw, 0.45, sc)
        txt(s, step, x=bx, y=4.3, w=sw, h=0.4,
            size=10.5, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        if i < len(flow_steps) - 1:
            txt(s, "→", x=bx + sw, y=4.32, w=0.18, h=0.35,
                size=12, color=arrow_color, align=PP_ALIGN.CENTER)

    return s


def slide_08_moscow(prs):
    """우선순위 MoSCoW"""
    s = blank_slide(prs)
    bg(s, C_LIGHT_BG)
    section_header(s, "우선순위 (MoSCoW)", y_pos=0.0)

    categories = {
        "Must": {
            "color": C_PRIMARY,
            "items": [
                "반려동물 프로필 등록",
                "데일리 체중 기록",
                "데일리 음수량 기록",
                "월간 건강 리포트",
            ],
        },
        "Should": {
            "color": C_ACCENT,
            "items": [
                "기록 리마인더 알림",
                "이상 징후 알림",
                "PDF 내보내기",
            ],
        },
        "Could": {
            "color": C_MUTED,
            "items": [
                "다두 펫 지원",
                "동물병원 공유 기능",
            ],
        },
    }

    col_w = 2.9
    for i, (cat, data) in enumerate(categories.items()):
        bx = 0.35 + i * 3.1
        col_h = 0.6 + len(data["items"]) * 0.65
        rect(s, bx, 0.68, col_w, col_h, C_CARD_BG)
        rect(s, bx, 0.68, col_w, 0.5, data["color"])
        txt(s, cat, x=bx + 0.1, y=0.7, w=col_w - 0.2, h=0.44,
            size=18, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        for j, item in enumerate(data["items"]):
            rect(s, bx + 0.12, 1.28 + j * 0.62, col_w - 0.24, 0.5, C_LIGHT_BG)
            txt(s, item,
                x=bx + 0.25, y=1.3 + j * 0.62, w=col_w - 0.4, h=0.44,
                size=13, color=C_DARK)

    # MVP 범위 안내
    rect(s, 0.3, 4.85, 9.4, 0.55, RGBColor(0xE0, 0xF5, 0xF7))
    txt(s, "MVP 범위: Must 전체 포함 | Should는 단계적 적용 | Could는 v2 이후 검토",
        x=0.4, y=4.9, w=9.2, h=0.42,
        size=13, bold=True, color=C_TAG_TEXT, align=PP_ALIGN.CENTER)

    return s


def slide_09_kpi(prs):
    """성공 지표 (KPI)"""
    s = blank_slide(prs)
    bg(s, C_LIGHT_BG)
    section_header(s, "성공 지표 (KPI) — 런칭 후 3개월 기준")

    kpis = [
        ("MAU",          "5,000",  "월간 활성 사용자"),
        ("7일 리텐션",   "≥ 40%",  "가입 후 1주일 재방문율"),
        ("기록 완료율",  "≥ 60%",  "활성 유저 데일리 기록"),
        ("리포트 열람율","≥ 70%",  "월간 리포트 오픈율"),
        ("평점",         "≥ 4.2",  "앱스토어 / 플레이스토어"),
    ]

    box_w = 1.72
    for i, (label, value, desc) in enumerate(kpis):
        bx = 0.3 + i * 1.88
        rect(s, bx, 0.75, box_w, 2.5, C_CARD_BG)
        rect(s, bx, 0.75, box_w, 0.06, C_ACCENT)
        txt(s, label,
            x=bx + 0.05, y=0.85, w=box_w - 0.1, h=0.35,
            size=12, bold=True, color=C_PRIMARY, align=PP_ALIGN.CENTER)
        txt(s, value,
            x=bx + 0.05, y=1.25, w=box_w - 0.1, h=0.7,
            size=30, bold=True, color=C_PRIMARY, align=PP_ALIGN.CENTER)
        txt(s, desc,
            x=bx + 0.05, y=2.0, w=box_w - 0.1, h=0.6,
            size=11, color=C_MUTED, align=PP_ALIGN.CENTER)

    # 설명 박스
    rect(s, 0.3, 3.5, 9.4, 1.7, C_CARD_BG)
    txt(s, "측정 전략",
        x=0.5, y=3.55, w=9.0, h=0.32,
        size=14, bold=True, color=C_PRIMARY)
    metrics = [
        "MAU: Supabase Auth 기반 월간 로그인 사용자 집계",
        "7일 리텐션: 가입일 +7일 이내 재방문 세션 여부 추적",
        "기록 완료율: 활성 유저(최근 7일 로그인) 중 일일 기록 완료 비율",
        "리포트 열람율: 월간 리포트 페이지 뷰 / 활성 유저 수",
    ]
    for i, m in enumerate(metrics):
        txt(s, f"· {m}",
            x=0.5, y=3.92 + i * 0.3, w=9.1, h=0.28,
            size=12, color=C_DARK)

    return s


def slide_10_timeline(prs):
    """출시 일정"""
    s = blank_slide(prs)
    bg(s, C_LIGHT_BG)
    section_header(s, "출시 일정 및 마일스톤")

    milestones = [
        ("M1", "프로젝트 착수",         "2026-05-19", "1.5주", C_MUTED),
        ("M2", "기획·디자인 확정",       "2026-06-12", "1.5주", C_PRIMARY),
        ("M3", "MVP 개발\n(체중·음수량·리포트)", "2026-07-05", "2.5주", C_ACCENT),
        ("M4", "QA & 베타 테스트",       "2026-07-17", "1.5주", C_PRIMARY),
        ("M5", "프로덕션 배포 출시",     "2026-07-22", "0.5주", C_ACCENT),
    ]

    total_weeks = 9.0
    start_x = 0.4
    bar_area_w = 9.2

    for i, (m_id, title, date, dur, color) in enumerate(milestones):
        y = 0.75 + i * 0.88
        weeks = float(dur.replace("주", ""))
        bar_w = (weeks / total_weeks) * bar_area_w

        # 전체 배경 바
        rect(s, start_x, y, bar_area_w, 0.65, C_CARD_BG)
        # 진행 바
        rect(s, start_x, y, bar_w, 0.65, color)

        # ID 배지
        txt(s, m_id,
            x=start_x + 0.08, y=y + 0.1, w=0.45, h=0.42,
            size=12, bold=True, color=C_WHITE)

        # 제목
        title_short = title.replace("\n", " ")
        txt(s, title_short,
            x=start_x + 0.6, y=y + 0.1, w=bar_w - 0.7, h=0.42,
            size=12.5, bold=True, color=C_WHITE)

        # 날짜 & 기간 (오른쪽)
        txt(s, f"{date}   ({dur})",
            x=start_x + bar_area_w - 2.2, y=y + 0.12, w=2.1, h=0.38,
            size=11, color=C_DARK if bar_w < 7.5 else C_WHITE,
            align=PP_ALIGN.RIGHT)

    # 현재 날짜 표시 마커
    today = "2026-05-28"
    progress_ratio = 9 / (9 * 7) * 7  # 9일 경과
    marker_x = start_x + (9 / (total_weeks * 7)) * bar_area_w
    line = s.shapes.add_shape(20,
                               Inches(start_x + 0.5), Inches(0.72),
                               Emu(0), Inches(4.5))
    line.line.color.rgb = RGBColor(0xFF, 0x66, 0x44)
    line.line.width = Pt(1.5)
    txt(s, f"오늘 {today}",
        x=start_x + 0.05, y=5.22, w=2.0, h=0.25,
        size=10, color=RGBColor(0xFF, 0x66, 0x44), bold=True)

    return s


def slide_11_open_items(prs):
    """미결 사항"""
    s = blank_slide(prs)
    bg(s, C_LIGHT_BG)
    section_header(s, "미결 사항 (Open Items)")

    items = [
        ("수익 모델",        "광고 기반 vs 구독 기반 결정 필요",
         "비즈니스 모델 미확정 — 기획 단계에서 방향 결정"),
        ("수의사 계정 (B2B)", "B2B 기능 범위 결정 필요",
         "수의사 전용 계정 및 처방 이력 공유 기능 검토"),
        ("추가 기록 항목",   "식사량, 배변 등 추가 여부 검토",
         "MVP 이후 단계에서 v2 요구사항으로 관리"),
        ("소셜 로그인",      "카카오·애플·구글 제공 범위 확정",
         "출시 전 법적 검토 및 UX 흐름 확정 필요"),
    ]

    for i, (title, summary, detail) in enumerate(items):
        bx = 0.3 + (i % 2) * 4.85
        by = 0.75 + (i // 2) * 2.2
        rect(s, bx, by, 4.5, 2.0, C_CARD_BG)
        rect(s, bx, by, 0.08, 2.0, C_MUTED)
        txt(s, title,
            x=bx + 0.22, y=by + 0.1, w=4.15, h=0.38,
            size=14, bold=True, color=C_PRIMARY)
        txt(s, summary,
            x=bx + 0.22, y=by + 0.52, w=4.15, h=0.35,
            size=12.5, bold=True, color=C_DARK)
        txt(s, detail,
            x=bx + 0.22, y=by + 0.9, w=4.15, h=0.85,
            size=12, color=C_MUTED)

    # 하단 액션 아이템
    rect(s, 0.3, 5.1, 9.4, 0.38, C_TAG_BG)
    txt(s, "액션: 수익 모델과 소셜 로그인 범위를 M2(기획 확정) 전까지 결정합니다.",
        x=0.45, y=5.13, w=9.1, h=0.3,
        size=12, bold=True, color=C_TAG_TEXT)

    return s


def slide_12_closing(prs):
    """마감 슬라이드"""
    s = blank_slide(prs)
    bg(s, C_PRIMARY)

    # 장식 원
    from pptx.util import Inches as I2
    circ = s.shapes.add_shape(9, I2(-0.5), I2(-0.5), I2(3), I2(3))
    circ.fill.solid()
    circ.fill.fore_color.rgb = RGBColor(0x00, 0x6B, 0x78)
    circ.line.fill.background()

    circ2 = s.shapes.add_shape(9, I2(7.5), I2(3.5), I2(3), I2(3))
    circ2.fill.solid()
    circ2.fill.fore_color.rgb = RGBColor(0x02, 0xA0, 0xB4)
    circ2.line.fill.background()

    # 세로 민트 바
    rect(s, 9.5, 0, 0.5, 5.625, C_ACCENT)

    txt(s, "감사합니다",
        x=1.5, y=1.4, w=7.0, h=1.0,
        size=44, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    line = s.shapes.add_shape(20, Inches(3), Inches(2.55),
                               Inches(4.0), Emu(0))
    line.line.color.rgb = C_ACCENT
    line.line.width = Pt(2)

    txt(s, "반려동물 건강 기록 웹앱",
        x=1.5, y=2.75, w=7.0, h=0.45,
        size=18, color=C_ACCENT, align=PP_ALIGN.CENTER)

    contact_lines = [
        "기간: 2026-05-19 ~ 2026-07-22 (9주)",
        "스택: React + Vite / Express.js / PostgreSQL (Supabase) / Vercel",
    ]
    for i, line_text in enumerate(contact_lines):
        txt(s, line_text,
            x=1.5, y=3.35 + i * 0.38, w=7.0, h=0.34,
            size=13, color=RGBColor(0xB2, 0xE4, 0xEA),
            align=PP_ALIGN.CENTER)

    return s


# ───────────────────────────────────────────────
# 메인 실행
# ───────────────────────────────────────────────

def main():
    prs = new_prs()

    print("슬라이드 생성 중...")
    slide_01_title(prs)                        ; print("  01/12 타이틀")
    slide_02_overview(prs)                     ; print("  02/12 프로젝트 개요")
    slide_03_problem(prs)                      ; print("  03/12 문제 정의")
    slide_04_users(prs)                        ; print("  04/12 대상 사용자")
    slide_05_feature_profile_weight(prs)       ; print("  05/12 핵심 기능(1)")
    slide_06_feature_water_report(prs)         ; print("  06/12 핵심 기능(2)")
    slide_07_feature_alerts(prs)               ; print("  07/12 핵심 기능(3)")
    slide_08_moscow(prs)                       ; print("  08/12 MoSCoW")
    slide_09_kpi(prs)                          ; print("  09/12 KPI")
    slide_10_timeline(prs)                     ; print("  10/12 출시 일정")
    slide_11_open_items(prs)                   ; print("  11/12 미결 사항")
    slide_12_closing(prs)                      ; print("  12/12 마감")

    timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M")
    output_path = f"05.리포트/PRD_발표자료_{timestamp}.pptx"
    prs.save(output_path)
    print(f"\n저장 완료: {output_path}")
    return output_path


if __name__ == "__main__":
    main()
